"""Office document tools: DOCX, XLSX, PPTX creation with ThinxS design system."""

from __future__ import annotations
import json
from pathlib import Path

from tools.sandbox import Sandbox
from tools.registry import ToolDef, TOOL_REGISTRY

# --- ThinxS Design System ---
# Shared palette and typography for all document types
THEME = {
    'primary': '0a1f1e',       # Deep teal (backgrounds, headers)
    'primary_light': '134e4a', # Lighter teal
    'accent': '14b8a6',        # Bright teal (highlights, accents)
    'accent_warm': 'f59e0b',   # Amber (callouts, emphasis)
    'text_light': 'f0fdfa',    # Light text on dark backgrounds
    'text_dark': '1a1a2e',     # Dark text on light backgrounds
    'surface': 'f5f0e8',       # Warm parchment (light surfaces)
    'surface_alt': 'ede8df',   # Alternate surface
    'border': '99f6e4',        # Teal border/accent lines
    'success': '4ade80',       # Green for status
    'muted': '6b7280',         # Gray for secondary text
    'white': 'ffffff',
    'font_heading': 'Segoe UI',
    'font_body': 'Segoe UI',
}


def register(sandbox: Sandbox):
    """Register document creation tools with the global registry."""

    def _output_path(filename: str) -> Path:
        """Resolve output path inside sandbox, creating output/ dir if needed."""
        out_dir = sandbox.allowed_dirs[0] / 'output'
        out_dir.mkdir(exist_ok=True)
        resolved = (out_dir / filename).resolve()
        sandbox.check_path(str(resolved))
        return resolved

    async def create_docx(filename: str, title: str, content: str,
                          author: str = 'Nemo Assistant') -> dict:
        """Create a Word document from structured content.

        content is a JSON string: a list of section objects.
        Each section: {"heading": "...", "body": "...", "level": 1}
        Or a plain string for simple paragraph text.
        """
        from docx import Document
        from docx.shared import Pt, Inches

        if not filename.endswith('.docx'):
            filename += '.docx'
        out = _output_path(filename)

        doc = Document()
        doc.core_properties.author = author
        doc.core_properties.title = title

        # Apply ThinxS theme to default styles
        from docx.shared import RGBColor
        style = doc.styles['Normal']
        font = style.font
        font.name = THEME['font_body']
        font.size = Pt(11)
        font.color.rgb = RGBColor.from_string(THEME['text_dark'])

        for level in range(4):
            heading_style = doc.styles[f'Heading {level + 1}'] if level > 0 else doc.styles['Title']
            hfont = heading_style.font
            hfont.name = THEME['font_heading']
            hfont.color.rgb = RGBColor.from_string(THEME['primary'])
            hfont.bold = True

        # Title
        doc.add_heading(title, level=0)

        # Parse content
        try:
            sections = json.loads(content)
        except (json.JSONDecodeError, TypeError):
            # Plain text fallback: split on double newlines for paragraphs
            sections = content

        if isinstance(sections, str):
            for para in sections.split('\n\n'):
                para = para.strip()
                if para:
                    doc.add_paragraph(para)
        elif isinstance(sections, list):
            for section in sections:
                if isinstance(section, str):
                    doc.add_paragraph(section)
                elif isinstance(section, dict):
                    if 'heading' in section:
                        level = section.get('level', 1)
                        doc.add_heading(section['heading'], level=min(level, 4))
                    if 'body' in section:
                        for para in section['body'].split('\n\n'):
                            para = para.strip()
                            if para:
                                doc.add_paragraph(para)
                    if 'bullets' in section:
                        for bullet in section['bullets']:
                            doc.add_paragraph(bullet, style='List Bullet')
                    if 'numbered' in section:
                        for item in section['numbered']:
                            doc.add_paragraph(item, style='List Number')
                    if 'table' in section:
                        table_data = section['table']
                        if table_data and len(table_data) > 0:
                            cols = len(table_data[0])
                            table = doc.add_table(rows=len(table_data), cols=cols)
                            table.style = 'Table Grid'
                            for r, row in enumerate(table_data):
                                for c, cell in enumerate(row):
                                    table.rows[r].cells[c].text = str(cell)

        doc.save(str(out))
        return {'message': f'Created {filename}', 'path': str(out), 'success': True}

    async def create_xlsx(filename: str, title: str, content: str) -> dict:
        """Create an Excel spreadsheet.

        content is a JSON string: either a list of lists (rows),
        or a dict of {"sheet_name": [[rows]]} for multiple sheets.
        First row of each sheet is treated as headers.
        """
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        if not filename.endswith('.xlsx'):
            filename += '.xlsx'
        out = _output_path(filename)

        wb = Workbook()
        wb.properties.title = title

        try:
            data = json.loads(content)
        except (json.JSONDecodeError, TypeError):
            return {'error': 'content must be valid JSON: list of lists or dict of sheet_name: [[rows]]'}

        header_font = Font(bold=True, color='FFFFFF', name=THEME['font_heading'])
        header_fill = PatternFill(start_color=THEME['primary'], end_color=THEME['primary'], fill_type='solid')
        alt_fill = PatternFill(start_color=THEME['surface'], end_color=THEME['surface'], fill_type='solid')
        thin_border = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin'),
        )

        def _populate_sheet(ws, rows):
            for r, row in enumerate(rows, 1):
                for c, val in enumerate(row, 1):
                    cell = ws.cell(row=r, column=c, value=val)
                    cell.border = thin_border
                    cell.font = Font(name=THEME['font_body'])
                    if r == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center')
                    elif r % 2 == 0:
                        cell.fill = alt_fill
            # Auto-width columns
            for c in range(1, ws.max_column + 1):
                max_len = 0
                for r in range(1, ws.max_row + 1):
                    val = ws.cell(row=r, column=c).value
                    if val:
                        max_len = max(max_len, len(str(val)))
                ws.column_dimensions[get_column_letter(c)].width = min(max_len + 4, 50)

        if isinstance(data, list):
            ws = wb.active
            ws.title = title[:31]  # Excel sheet name limit
            _populate_sheet(ws, data)
        elif isinstance(data, dict):
            first = True
            for sheet_name, rows in data.items():
                if first:
                    ws = wb.active
                    ws.title = sheet_name[:31]
                    first = False
                else:
                    ws = wb.create_sheet(title=sheet_name[:31])
                _populate_sheet(ws, rows)
        else:
            return {'error': 'content must be a list of lists or dict of sheets'}

        wb.save(str(out))
        return {'message': f'Created {filename}', 'path': str(out), 'success': True}

    async def create_pptx(filename: str, title: str, content: str,
                          author: str = 'Nemo Assistant') -> dict:
        """Create a styled PowerPoint presentation with ThinxS design system.

        content is a JSON string: a list of slide objects.
        Each slide: {"title": "...", "body": "...", "bullets": [...],
                     "layout": "title|section|content|two_column|blank",
                     "notes": "...", "subtitle": "...",
                     "left_bullets": [...], "right_bullets": [...],
                     "left_heading": "...", "right_heading": "..."}
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor

        if not filename.endswith('.pptx'):
            filename += '.pptx'
        out = _output_path(filename)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        prs.core_properties.author = author
        prs.core_properties.title = title

        try:
            slides = json.loads(content)
        except (json.JSONDecodeError, TypeError):
            return {'error': 'content must be valid JSON: list of slide objects'}

        if not isinstance(slides, list):
            return {'error': 'content must be a JSON list of slide objects'}

        T = THEME  # shorthand

        def _rgb(hex_str):
            return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

        def _set_slide_bg(slide, color_hex):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(color_hex)

        def _add_shape(slide, left, top, width, height, fill_hex=None):
            from pptx.util import Emu
            shape = slide.shapes.add_shape(
                1, left, top, width, height  # 1 = rectangle
            )
            shape.line.fill.background()  # no border
            if fill_hex:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(fill_hex)
            return shape

        def _style_title(tf, text, font_size=36, color=None, bold=True, align=PP_ALIGN.LEFT):
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = T['font_heading']
            if color:
                run.font.color.rgb = _rgb(color)

        def _add_text_box(slide, left, top, width, height, text, font_size=18,
                          color=None, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            try:
                tf.vertical_anchor = anchor
            except Exception:
                pass
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = T['font_body']
            run.font.bold = bold
            if color:
                run.font.color.rgb = _rgb(color)
            return tf

        def _add_bullets(tf, items, font_size=16, color=None, start_new=False):
            for i, item in enumerate(items):
                if i == 0 and not start_new and not tf.paragraphs[0].text:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.space_after = Pt(6)
                p.space_before = Pt(2)
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(font_size)
                run.font.name = T['font_body']
                if color:
                    run.font.color.rgb = _rgb(color)

        def _add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.8)):
            bar = _add_shape(slide, left, top, width, height, T['accent'])
            return bar

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for i, slide_data in enumerate(slides):
            if isinstance(slide_data, str):
                slide_data = {'title': slide_data}

            layout_name = slide_data.get('layout', 'content')
            slide_title = slide_data.get('title', '')
            subtitle = slide_data.get('subtitle', '')
            body = slide_data.get('body', '')
            bullets = slide_data.get('bullets', [])
            notes = slide_data.get('notes', '')

            # Use blank layout for all slides (we draw everything custom)
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            if layout_name == 'title' or (i == 0 and layout_name != 'blank'):
                # --- TITLE SLIDE: dark bg, centered title, accent line ---
                _set_slide_bg(slide, T['primary'])

                # Accent line across top
                _add_shape(slide, Inches(0), Inches(0), slide_w, Inches(0.06), T['accent'])

                # Title centered
                _add_text_box(
                    slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(2),
                    slide_title, font_size=44, color=T['text_light'],
                    bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
                )

                # Accent divider
                _add_shape(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.05), T['accent'])

                # Subtitle
                if subtitle:
                    _add_text_box(
                        slide, Inches(2), Inches(4.6), Inches(9.3), Inches(1),
                        subtitle, font_size=20, color=T['border'],
                        align=PP_ALIGN.CENTER
                    )

                # Author/footer
                _add_text_box(
                    slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.6),
                    f'{author}', font_size=12, color=T['muted'],
                    align=PP_ALIGN.CENTER
                )

            elif layout_name == 'section':
                # --- SECTION DIVIDER: dark bg, large title ---
                _set_slide_bg(slide, T['primary_light'])
                _add_shape(slide, Inches(0), Inches(0), slide_w, Inches(0.06), T['accent'])

                _add_text_box(
                    slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2),
                    slide_title, font_size=40, color=T['text_light'],
                    bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE
                )

                if body:
                    _add_text_box(
                        slide, Inches(1.5), Inches(4.5), Inches(9), Inches(1.5),
                        body, font_size=18, color=T['border']
                    )

            elif layout_name == 'two_column':
                # --- TWO COLUMN: light bg, side by side content ---
                _set_slide_bg(slide, T['surface'])
                _add_shape(slide, Inches(0), Inches(0), slide_w, Inches(0.06), T['accent'])

                # Title bar
                _add_shape(slide, Inches(0), Inches(0.06), slide_w, Inches(1.2), T['primary'])
                _add_text_box(
                    slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(1),
                    slide_title, font_size=30, color=T['text_light'], bold=True,
                    anchor=MSO_ANCHOR.MIDDLE
                )

                left_heading = slide_data.get('left_heading', '')
                right_heading = slide_data.get('right_heading', '')
                left_bullets = slide_data.get('left_bullets', [])
                right_bullets = slide_data.get('right_bullets', [])

                col_top = Inches(1.6)
                col_width = Inches(5.5)

                # Left column
                if left_heading:
                    _add_accent_bar(slide, Inches(0.7), col_top, height=Inches(0.5))
                    _add_text_box(
                        slide, Inches(0.95), col_top, col_width, Inches(0.5),
                        left_heading, font_size=20, color=T['primary'], bold=True
                    )
                if left_bullets:
                    tf = _add_text_box(
                        slide, Inches(0.95), col_top + Inches(0.7), col_width, Inches(4.5),
                        '', font_size=16, color=T['text_dark']
                    )
                    _add_bullets(tf, left_bullets, color=T['text_dark'])

                # Right column
                if right_heading:
                    _add_accent_bar(slide, Inches(6.9), col_top, height=Inches(0.5))
                    _add_text_box(
                        slide, Inches(7.15), col_top, col_width, Inches(0.5),
                        right_heading, font_size=20, color=T['primary'], bold=True
                    )
                if right_bullets:
                    tf = _add_text_box(
                        slide, Inches(7.15), col_top + Inches(0.7), col_width, Inches(4.5),
                        '', font_size=16, color=T['text_dark']
                    )
                    _add_bullets(tf, right_bullets, color=T['text_dark'])

            elif layout_name == 'blank':
                _set_slide_bg(slide, T['surface'])

            else:
                # --- CONTENT SLIDE: light bg, title bar, body + bullets ---
                _set_slide_bg(slide, T['surface'])

                # Top accent line
                _add_shape(slide, Inches(0), Inches(0), slide_w, Inches(0.06), T['accent'])

                # Title bar
                _add_shape(slide, Inches(0), Inches(0.06), slide_w, Inches(1.2), T['primary'])
                _add_text_box(
                    slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(1),
                    slide_title, font_size=30, color=T['text_light'], bold=True,
                    anchor=MSO_ANCHOR.MIDDLE
                )

                content_top = Inches(1.6)

                # Accent bar next to content
                if bullets or body:
                    _add_accent_bar(slide, Inches(0.7), content_top)

                # Body text
                if body:
                    body_tf = _add_text_box(
                        slide, Inches(0.95), content_top, Inches(11), Inches(1.5),
                        body, font_size=18, color=T['text_dark']
                    )
                    content_top += Inches(1.5)

                # Bullets
                if bullets:
                    bullet_top = content_top + (Inches(0.2) if body else Inches(0))
                    tf = _add_text_box(
                        slide, Inches(0.95), bullet_top, Inches(11), Inches(4.5),
                        '', font_size=16, color=T['text_dark']
                    )
                    _add_bullets(tf, bullets, font_size=18, color=T['text_dark'])

            # Slide number (bottom right, all slides except title)
            if not (layout_name == 'title' or (i == 0 and layout_name != 'blank')):
                _add_text_box(
                    slide, Inches(12), Inches(7), Inches(1), Inches(0.4),
                    str(i + 1), font_size=10, color=T['muted'],
                    align=PP_ALIGN.RIGHT
                )

            # Speaker notes
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        prs.save(str(out))
        return {
            'message': f'Created {filename} with {len(slides)} slides',
            'path': str(out),
            'success': True,
        }

    # Register tools
    TOOL_REGISTRY.register(ToolDef(
        name='create_docx',
        description='Create a Word document (.docx). Content can be plain text or JSON array of sections with headings, body, bullets, numbered lists, and tables.',
        parameters={
            'filename': {'type': 'string', 'description': 'Output filename (e.g. "report.docx")'},
            'title': {'type': 'string', 'description': 'Document title'},
            'content': {'type': 'string', 'description': 'Plain text or JSON array of section objects: [{"heading": "...", "body": "...", "bullets": [...], "table": [[...]]}]'},
            'author': {'type': 'string', 'description': 'Author name (default: Nemo Assistant)'},
        },
        required=['filename', 'title', 'content'],
        handler=create_docx,
    ))

    TOOL_REGISTRY.register(ToolDef(
        name='create_xlsx',
        description='Create an Excel spreadsheet (.xlsx). Content is JSON: a list of lists (rows) for single sheet, or {"sheet_name": [[rows]]} for multiple sheets. First row = headers.',
        parameters={
            'filename': {'type': 'string', 'description': 'Output filename (e.g. "data.xlsx")'},
            'title': {'type': 'string', 'description': 'Workbook title'},
            'content': {'type': 'string', 'description': 'JSON: list of lists [[row1], [row2]] or dict {"Sheet1": [[rows]], "Sheet2": [[rows]]}'},
        },
        handler=create_xlsx,
    ))

    TOOL_REGISTRY.register(ToolDef(
        name='create_pptx',
        description=(
            'Create a professionally styled PowerPoint presentation (.pptx) with ThinxS design system '
            '(dark teal headers, warm parchment backgrounds, accent bars). '
            'Layouts: "title" (dark bg, centered), "section" (divider slide), "content" (standard), '
            '"two_column" (side-by-side with left/right headings and bullets), "blank".'
        ),
        parameters={
            'filename': {'type': 'string', 'description': 'Output filename (e.g. "deck.pptx")'},
            'title': {'type': 'string', 'description': 'Presentation title'},
            'content': {'type': 'string', 'description': (
                'JSON list of slides: [{"title": "...", "subtitle": "...", "body": "...", '
                '"bullets": [...], "layout": "title|section|content|two_column|blank", '
                '"left_heading": "...", "left_bullets": [...], "right_heading": "...", '
                '"right_bullets": [...], "notes": "..."}]'
            )},
            'author': {'type': 'string', 'description': 'Author name (default: Nemo Assistant)'},
        },
        required=['filename', 'title', 'content'],
        handler=create_pptx,
    ))
